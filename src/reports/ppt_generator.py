"""
Gerador de Apresentações PowerPoint Automatizadas
Cria apresentações profissionais com gráficos e tabelas
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from datetime import datetime
import pandas as pd


class PowerPointGenerator:
    """
    Gerador de apresentações PowerPoint automatizadas.

    Cria apresentação com:
    - Slide título para cada cenário
    - Gráficos PNG embedados
    - Tabelas de resultados
    - Total: 15-20 slides
    """

    def __init__(self, output_dir: str = "reports/powerpoint"):
        """
        Inicializa o gerador de PowerPoint.

        Args:
            output_dir: Diretório para salvar as apresentações
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Cores corporativas
        self.cor_primaria = RGBColor(54, 96, 146)  # Azul escuro
        self.cor_secundaria = RGBColor(79, 129, 189)  # Azul claro
        self.cor_destaque = RGBColor(192, 0, 0)  # Vermelho
        self.cor_sucesso = RGBColor(0, 176, 80)  # Verde

    def _adicionar_slide_titulo(self, prs: Presentation, titulo: str, subtitulo: str = ""):
        """Adiciona slide de título"""
        slide_layout = prs.slide_layouts[0]  # Layout de título
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.cor_primaria

        if subtitulo:
            subtitle.text = subtitulo
            subtitle.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle.text_frame.paragraphs[0].font.color.rgb = self.cor_secundaria

        return slide

    def _adicionar_slide_secao(self, prs: Presentation, titulo: str):
        """Adiciona slide de seção"""
        slide_layout = prs.slide_layouts[2]  # Layout de seção
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Fundo colorido
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.cor_primaria

        return slide

    def _adicionar_slide_conteudo(
        self,
        prs: Presentation,
        titulo: str,
        conteudo: Optional[str] = None
    ):
        """Adiciona slide de conteúdo"""
        slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.cor_primaria

        if conteudo:
            body = slide.placeholders[1]
            body.text = conteudo
            body.text_frame.paragraphs[0].font.size = Pt(18)

        return slide

    def _adicionar_imagem(
        self,
        slide,
        caminho_imagem: Path,
        left: float = 1.5,
        top: float = 2.0,
        width: float = 7.0
    ):
        """Adiciona imagem ao slide"""
        if caminho_imagem.exists():
            slide.shapes.add_picture(
                str(caminho_imagem),
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
        else:
            print(f"⚠ Imagem não encontrada: {caminho_imagem}")

    def _adicionar_tabela(
        self,
        slide,
        df: pd.DataFrame,
        left: float = 1.5,
        top: float = 2.5,
        width: float = 7.0,
        height: float = 3.0
    ):
        """Adiciona tabela ao slide"""
        rows, cols = df.shape
        rows += 1  # Header

        # Cria tabela
        table = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        ).table

        # Preenche header
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.cor_primaria
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Preenche dados
        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Zebra striping
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

    def _adicionar_bullet_points(
        self,
        slide,
        titulo: str,
        pontos: List[str]
    ):
        """Adiciona slide com bullet points"""
        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.cor_primaria

        # Adiciona caixa de texto para bullets
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(7.5)
        height = Inches(4.0)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, ponto in enumerate(pontos):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = ponto
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(10)

    def slide_capa(
        self,
        prs: Presentation,
        titulo: str = "Análise de Viés em Avaliações de RH",
        subtitulo: Optional[str] = None
    ):
        """Slide 1: Capa"""
        if subtitulo is None:
            subtitulo = f"Relatório Gerado em {datetime.now().strftime('%d/%m/%Y')}"

        self._adicionar_slide_titulo(prs, titulo, subtitulo)
        print("✓ Slide 1: Capa")

    def slide_agenda(self, prs: Presentation, num_cenarios: int = 7):
        """Slide 2: Agenda"""
        slide = self._adicionar_slide_conteudo(prs, "Agenda")

        pontos = [
            "1. Introdução e Metodologia",
            f"2. Análise de {num_cenarios} Cenários de Correção de Viés",
            "3. Análise Comparativa",
            "4. Conclusões e Recomendações"
        ]

        self._adicionar_bullet_points(slide, "Agenda", pontos)
        print("✓ Slide 2: Agenda")

    def slide_metodologia(self, prs: Presentation):
        """Slide 3: Metodologia"""
        slide = self._adicionar_slide_conteudo(prs, "Metodologia")

        pontos = [
            "Framework de detecção e correção de viés de gênero",
            "Análise estatística com testes t de Student",
            "Nível de significância: α = 0.05",
            "Correção por reponderação de scores",
            "Avaliação de eficácia da correção"
        ]

        self._adicionar_bullet_points(slide, "Metodologia", pontos)
        print("✓ Slide 3: Metodologia")

    def slide_cenario_titulo(
        self,
        prs: Presentation,
        numero: int,
        titulo: str,
        descricao: str
    ):
        """Slide de título do cenário"""
        self._adicionar_slide_secao(prs, f"Cenário {numero}")

        # Adiciona slide com descrição
        slide = self._adicionar_slide_conteudo(prs, titulo, descricao)

        print(f"✓ Slide: Cenário {numero} - {titulo}")

    def slide_grafico(
        self,
        prs: Presentation,
        titulo: str,
        caminho_grafico: Path,
        observacoes: Optional[str] = None
    ):
        """Slide com gráfico"""
        slide = self._adicionar_slide_conteudo(prs, titulo)

        # Adiciona gráfico
        self._adicionar_imagem(slide, caminho_grafico, left=1.0, top=1.8, width=8.0)

        # Adiciona observações se houver
        if observacoes:
            left = Inches(1.0)
            top = Inches(5.8)
            width = Inches(8.0)
            height = Inches(0.8)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.text = observacoes
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = RGBColor(100, 100, 100)

        print(f"✓ Slide: {titulo}")

    def slide_tabela_resultados(
        self,
        prs: Presentation,
        titulo: str,
        df: pd.DataFrame
    ):
        """Slide com tabela de resultados"""
        slide = self._adicionar_slide_conteudo(prs, titulo)

        # Adiciona tabela
        self._adicionar_tabela(slide, df, left=1.0, top=2.0, width=8.0, height=3.5)

        print(f"✓ Slide: {titulo}")

    def slide_comparativo(
        self,
        prs: Presentation,
        titulo: str,
        dados_comparativo: Dict[str, float]
    ):
        """Slide comparativo entre cenários"""
        slide = self._adicionar_slide_conteudo(prs, titulo)

        # Cria DataFrame para tabela
        df = pd.DataFrame([dados_comparativo])

        self._adicionar_tabela(slide, df, left=1.0, top=2.5, width=8.0, height=2.0)

        print(f"✓ Slide: {titulo}")

    def slide_conclusoes(
        self,
        prs: Presentation,
        conclusoes: List[str]
    ):
        """Slide de conclusões"""
        slide = self._adicionar_slide_conteudo(prs, "Conclusões")

        self._adicionar_bullet_points(slide, "Conclusões", conclusoes)

        print("✓ Slide: Conclusões")

    def slide_recomendacoes(
        self,
        prs: Presentation,
        recomendacoes: List[str]
    ):
        """Slide de recomendações"""
        slide = self._adicionar_slide_conteudo(prs, "Recomendações")

        self._adicionar_bullet_points(slide, "Recomendações", recomendacoes)

        print("✓ Slide: Recomendações")

    def slide_final(self, prs: Presentation):
        """Slide final"""
        self._adicionar_slide_titulo(
            prs,
            "Obrigado!",
            "Framework de Redução de Viés em Avaliações de RH"
        )

        print("✓ Slide: Final")

    def gerar_apresentacao_completa(
        self,
        graficos: List[Path],
        tabelas: Dict[str, pd.DataFrame],
        dados_cenarios: Dict,
        nome_arquivo: Optional[str] = None
    ) -> Path:
        """
        Gera apresentação PowerPoint completa com N cenários.

        Args:
            graficos: Lista de caminhos dos gráficos PNG
            tabelas: Dicionário com DataFrames para tabelas (chaves: 'cenario_1', 'cenario_2', etc)
            dados_cenarios: Dados dos cenários (qualquer número)
                {
                    'cenario_1': {'titulo': '...', 'descricao': '...'},
                    'cenario_2': {...},
                    ...
                }
            nome_arquivo: Nome do arquivo (opcional)

        Returns:
            Path do arquivo gerado
        """
        if nome_arquivo is None:
            nome_arquivo = f"apresentacao_vies_{self.timestamp}.pptx"

        caminho = self.output_dir / nome_arquivo

        print("\n=== Gerando Apresentação PowerPoint ===\n")

        # Cria apresentação
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Capa
        self.slide_capa(prs)

        # Slide 2: Agenda
        self.slide_agenda(prs, num_cenarios=len(dados_cenarios))

        # Slide 3: Metodologia
        self.slide_metodologia(prs)

        # Itera pelos cenários dinamicamente
        idx_grafico = 0
        for key in sorted(dados_cenarios.keys()):
            # Extrai número do cenário
            num_cenario = int(key.split('_')[-1]) if '_' in key else 0
            dados = dados_cenarios[key]

            # Slide de título do cenário
            titulo = dados.get('titulo', f'Cenário {num_cenario}')
            descricao = dados.get('descricao', f'Análise do {titulo}')

            self.slide_cenario_titulo(
                prs,
                num_cenario,
                titulo,
                descricao
            )

            # Adiciona gráfico se disponível
            if idx_grafico < len(graficos) and graficos[idx_grafico].exists():
                self.slide_grafico(
                    prs,
                    f"Distribuição de Scores - {titulo}",
                    graficos[idx_grafico],
                    dados.get('observacao_grafico', '')
                )
                idx_grafico += 1

            # Adiciona tabela se disponível
            if key in tabelas:
                self.slide_tabela_resultados(
                    prs,
                    f"Resultados Estatísticos - {titulo}",
                    tabelas[key]
                )

        # Análise Comparativa
        self._adicionar_slide_secao(prs, "Análise Comparativa")

        if len(graficos) > 3:
            self.slide_grafico(
                prs,
                "Comparação entre Cenários",
                graficos[3]
            )

        # Gráficos adicionais
        graficos_extras = [
            ("Eficácia da Correção", 4),
            ("Distribuição Geral", 5),
            ("Desempenho vs Potencial", 6)
        ]

        for titulo, idx in graficos_extras:
            if len(graficos) > idx:
                self.slide_grafico(prs, titulo, graficos[idx])

        # Conclusões
        conclusoes = [
            "Viés de gênero foi detectado nos dados originais",
            "A correção por reponderação mostrou-se eficaz",
            "Redução significativa na diferença de médias entre gêneros",
            "P-values indicam melhoria na equidade estatística",
            "Mudanças no ranking refletem correção aplicada"
        ]
        self.slide_conclusoes(prs, conclusoes)

        # Recomendações
        recomendacoes = [
            "Implementar correção de viés como prática padrão",
            "Monitorar métricas de equidade periodicamente",
            "Treinar avaliadores sobre vieses inconscientes",
            "Revisar processos de avaliação regularmente",
            "Documentar decisões e justificativas"
        ]
        self.slide_recomendacoes(prs, recomendacoes)

        # Slide final
        self.slide_final(prs)

        # Salva apresentação
        prs.save(str(caminho))

        print(f"\n✓ Apresentação PowerPoint salva: {caminho}")
        print(f"  Total de slides: {len(prs.slides)}")

        return caminho
