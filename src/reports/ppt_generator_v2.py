"""
Gerador de Apresentações PowerPoint Profissionais - VERSÃO MELHORADA
Cria apresentações de alto impacto com análises detalhadas
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from datetime import datetime
import pandas as pd


class PowerPointGeneratorV2:
    """
    Gerador de apresentações PowerPoint profissionais - VERSÃO MELHORADA

    Melhorias:
    - Design visual mais impactante
    - Análises mais detalhadas
    - Insights e recomendações por cenário
    - Gráficos maiores e mais claros
    - Narrativa coesa e profissional
    """

    def __init__(self, output_dir: str = "reports/powerpoint"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Paleta de cores profissional
        self.cor_primaria = RGBColor(26, 35, 126)  # Azul escuro profissional
        self.cor_secundaria = RGBColor(63, 81, 181)  # Azul médio
        self.cor_destaque = RGBColor(244, 67, 54)  # Vermelho para alertas
        self.cor_sucesso = RGBColor(76, 175, 80)  # Verde para sucesso
        self.cor_aviso = RGBColor(255, 152, 0)  # Laranja para avisos
        self.cor_texto = RGBColor(33, 33, 33)  # Cinza escuro
        self.cor_texto_claro = RGBColor(117, 117, 117)  # Cinza claro

    def _criar_caixa_destaque(self, slide, texto: str, left: float, top: float,
                             width: float, height: float, cor: RGBColor):
        """Cria uma caixa de destaque colorida"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        # Configurar preenchimento
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = cor

        # Configurar texto
        text_frame = shape.text_frame
        text_frame.text = texto
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        return shape

    def _adicionar_rodape(self, slide, texto: str = "Framework de Redução de Viés"):
        """Adiciona rodapé profissional"""
        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(7),
            Inches(9), Inches(0.3)
        )
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(10)
        p.font.color.rgb = self.cor_texto_claro
        p.alignment = PP_ALIGN.CENTER

    def slide_capa_melhorada(self, prs: Presentation):
        """Slide de capa profissional"""
        slide_layout = prs.slide_layouts[6]  # Layout em branco
        slide = prs.slides.add_slide(slide_layout)

        # Fundo gradiente simulado com retângulo
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.cor_primaria

        # Título principal
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(1.5)
        )
        text_frame = title_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Análise de Viés em Avaliações de RH"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtítulo
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8),
            Inches(8), Inches(0.8)
        )
        text_frame = subtitle_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Framework de Detecção e Correção de Viés de Gênero"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Análise de 7 cenários
        info_box = slide.shapes.add_textbox(
            Inches(2), Inches(5),
            Inches(6), Inches(0.6)
        )
        text_frame = info_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Análise Progressiva: 0% → 100% de Correção"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 235, 59)  # Amarelo
        p.alignment = PP_ALIGN.CENTER

        # Data
        date_box = slide.shapes.add_textbox(
            Inches(3), Inches(6.5),
            Inches(4), Inches(0.5)
        )
        text_frame = date_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = datetime.now().strftime("%B %Y")
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

        print("✓ Slide 1: Capa Profissional")

    def slide_contexto_problema(self, prs: Presentation):
        """Slide explicando o problema de viés"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "O Problema: Viés de Gênero em Avaliações"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cor_primaria

        # Caixas de problemas
        problemas = [
            ("Diferença de 0.56 pontos", "Entre avaliações de homens e mulheres"),
            ("P-value < 0.0001", "Estatisticamente significativo"),
            ("Impacto no ranking", "Afeta decisões de promoção e remuneração")
        ]

        for i, (titulo_prob, desc) in enumerate(problemas):
            top = 2.2 + (i * 1.5)

            # Caixa principal
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(top),
                Inches(7), Inches(1.2)
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.cor_destaque if i == 0 else self.cor_aviso

            # Título
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = titulo_prob
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Descrição
            p = text_frame.add_paragraph()
            p.text = desc
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(6)

        self._adicionar_rodape(slide)
        print("✓ Slide: Contexto do Problema")

    def slide_cenario_detalhado(
        self,
        prs: Presentation,
        numero: int,
        dados: Dict,
        imagem_path: Optional[Path] = None
    ):
        """Slide de cenário com análise detalhada"""
        slide_layout = prs.slide_layouts[6]  # Branco
        slide = prs.slides.add_slide(slide_layout)

        # Header colorido
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        fill = header.fill
        fill.solid()

        # Cor baseada no nível de correção
        nivel = dados.get('nivel_correcao', 0)
        if nivel == 0:
            fill.fore_color.rgb = self.cor_destaque
        elif nivel < 50:
            fill.fore_color.rgb = self.cor_aviso
        elif nivel < 80:
            fill.fore_color.rgb = RGBColor(255, 193, 7)
        else:
            fill.fore_color.rgb = self.cor_sucesso

        # Título do cenário
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(7), Inches(0.8)
        )
        text_frame = title_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = dados.get('titulo', f'Cenário {numero}')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Nível de correção
        nivel_box = slide.shapes.add_textbox(
            Inches(7.5), Inches(0.3),
            Inches(2), Inches(0.6)
        )
        text_frame = nivel_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"{nivel:.0f}% Correção"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT

        # Métricas principais em cards
        metricas = [
            ("Diferença", f"{abs(dados.get('diferenca_depois', 0)):.3f}",
             "Entre médias" ),
            ("P-value", f"{dados.get('p_value_depois', 0):.4f}",
             "Significância"),
            ("Status", "✓ Sem Viés" if dados.get('p_value_depois', 1) >= 0.05 else "⚠ Com Viés",
             "Resultado")
        ]

        for i, (label, valor, sublabel) in enumerate(metricas):
            left = 0.5 + (i * 3.2)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(1.5),
                Inches(3), Inches(1.2)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 245, 245)

            text_frame = card.text_frame
            text_frame.clear()

            # Label
            p = text_frame.add_paragraph()
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.cor_texto_claro
            p.alignment = PP_ALIGN.CENTER

            # Valor
            p = text_frame.add_paragraph()
            p.text = valor
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.cor_primaria
            p.alignment = PP_ALIGN.CENTER

            # Sublabel
            p = text_frame.add_paragraph()
            p.text = sublabel
            p.font.size = Pt(12)
            p.font.color.rgb = self.cor_texto_claro
            p.alignment = PP_ALIGN.CENTER

        # Gráfico (se disponível)
        if imagem_path and imagem_path.exists():
            slide.shapes.add_picture(
                str(imagem_path),
                Inches(0.5), Inches(3),
                width=Inches(9)
            )

        # Insights
        insights_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.3),
            Inches(9), Inches(0.8)
        )
        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = f"💡 Insight: {dados.get('insight', 'Análise em andamento...')}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self.cor_texto

        self._adicionar_rodape(slide)
        print(f"✓ Slide: {dados.get('titulo', f'Cenário {numero}')}")

    def slide_comparativo_evolucao(
        self,
        prs: Presentation,
        dados_cenarios: Dict
    ):
        """Slide mostrando evolução entre cenários"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.6)
        )
        text_frame = title_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Evolução da Correção de Viés"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.cor_primaria
        p.alignment = PP_ALIGN.CENTER

        # Tabela de evolução
        rows = len(dados_cenarios) + 1
        cols = 5

        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.8), Inches(1.5),
            Inches(8.4), Inches(4.5)
        ).table

        # Header
        headers = ["Cenário", "Correção", "Diferença", "P-value", "Status"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.cor_primaria

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Dados
        for row_idx, (key, dados) in enumerate(sorted(dados_cenarios.items()), 1):
            # Cenário
            cell = table.cell(row_idx, 0)
            cell.text = f"Cenário {row_idx}"
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # Correção
            nivel = dados.get('nivel_correcao', 0)
            cell = table.cell(row_idx, 1)
            cell.text = f"{nivel:.0f}%"
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Diferença
            diff = abs(dados.get('diferenca_depois', 0))
            cell = table.cell(row_idx, 2)
            cell.text = f"{diff:.3f}"
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # P-value
            pval = dados.get('p_value_depois', 1)
            cell = table.cell(row_idx, 3)
            cell.text = f"{pval:.4f}"
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # Status com cor
            cell = table.cell(row_idx, 4)
            if pval < 0.05:
                cell.text = "⚠ Com Viés"
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 235, 238)
            else:
                cell.text = "✓ Sem Viés"
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 245, 233)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

        self._adicionar_rodape(slide)
        print("✓ Slide: Comparativo de Evolução")

    def gerar_apresentacao_profissional(
        self,
        dados_cenarios: Dict,
        graficos: List[Path] = None,
        nome_arquivo: Optional[str] = None
    ) -> Path:
        """Gera apresentação PowerPoint profissional"""
        if nome_arquivo is None:
            nome_arquivo = f"apresentacao_profissional_{self.timestamp}.pptx"

        caminho = self.output_dir / nome_arquivo

        print("\n=== Gerando Apresentação PowerPoint PROFISSIONAL ===\n")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Capa
        self.slide_capa_melhorada(prs)

        # Slide 2: Contexto do problema
        self.slide_contexto_problema(prs)

        # Slides dos cenários com análise detalhada
        for idx, (key, dados) in enumerate(sorted(dados_cenarios.items()), 1):
            # Adiciona nível de correção aos dados
            niveis = [0, 16.67, 33.33, 50, 66.67, 83.33, 100]
            dados['nivel_correcao'] = niveis[idx - 1] if idx <= len(niveis) else 0

            # Gera insights baseados nos dados
            pval = dados.get('p_value_depois', 1)
            diff = abs(dados.get('diferenca_depois', 0))

            if pval < 0.001:
                dados['insight'] = f"Viés muito significativo detectado. Diferença de {diff:.3f} pontos indica desigualdade sistêmica."
            elif pval < 0.05:
                dados['insight'] = f"Viés estatisticamente significativo. A diferença de {diff:.3f} ainda requer atenção."
            elif pval < 0.10:
                dados['insight'] = f"Viés marginal. Diferença de {diff:.3f} está no limite de significância."
            else:
                dados['insight'] = f"Equidade alcançada! Diferença de {diff:.3f} não é estatisticamente significativa."

            # Imagem do cenário (se disponível)
            img_path = graficos[idx-1] if graficos and idx <= len(graficos) else None

            self.slide_cenario_detalhado(prs, idx, dados, img_path)

        # Slide comparativo
        self.slide_comparativo_evolucao(prs, dados_cenarios)

        # Salvar
        prs.save(str(caminho))

        print(f"\n✓ Apresentação PROFISSIONAL salva: {caminho}")
        print(f"  Total de slides: {len(prs.slides)}")

        return caminho
