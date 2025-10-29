"""
Gerador de Apresentações PowerPoint ULTRA DETALHADAS - V3
Cada cenário tem múltiplos slides com TODOS os gráficos e análises completas
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime


class PowerPointGeneratorV3:
    """
    Gerador PowerPoint ULTRA DETALHADO

    Cada cenário recebe:
    - 1 slide de introdução com métricas
    - 8 slides com gráficos (1 por gráfico)
    - 1 slide de conclusão
    Total: ~70 slides completos
    """

    def __init__(self, output_dir: str = "reports/powerpoint"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Cores
        self.cor_sem_vies = RGBColor(76, 175, 80)  # Verde
        self.cor_vies_leve = RGBColor(255, 193, 7)  # Amarelo
        self.cor_vies_moderado = RGBColor(255, 152, 0)  # Laranja
        self.cor_vies_severo = RGBColor(244, 67, 54)  # Vermelho

    def _get_cor_cenario(self, nivel_vies: float) -> RGBColor:
        """Retorna cor baseada no nível de viés"""
        if nivel_vies == 0:
            return self.cor_sem_vies
        elif nivel_vies < 50:
            return self.cor_vies_leve
        elif nivel_vies < 80:
            return self.cor_vies_moderado
        else:
            return self.cor_vies_severo

    def _slide_capa(self, prs: Presentation):
        """Slide de capa"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Fundo
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(26, 35, 126)

        # Título
        title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Análise Completa de Viés\nem Avaliações de RH"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtítulo
        sub = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(0.8))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = "7 Cenários Completos | 56 Gráficos | Análises Detalhadas"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 235, 59)
        p.alignment = PP_ALIGN.CENTER

        print("✓ Slide: Capa")

    def _slide_indice(self, prs: Presentation, dados_cenarios: Dict):
        """Slide de índice"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Índice de Cenários"

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for key in sorted(dados_cenarios.keys()):
            dados = dados_cenarios[key]
            num = int(key.split('_')[1])

            p = tf.add_paragraph()
            p.text = f"Cenário {num}: {dados['titulo']}"
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(10)

            # Sub-item com status
            p = tf.add_paragraph()
            status = "✅ SEM VIÉS" if not dados['vies_detectado'] else "⚠️ VIÉS DETECTADO"
            p.text = f"{status} | Diferença: {abs(dados['diferenca_depois']):.3f}"
            p.level = 1
            p.font.size = Pt(14)

        print("✓ Slide: Índice")

    def _slide_intro_cenario(
        self,
        prs: Presentation,
        numero: int,
        dados: Dict
    ):
        """Slide de introdução do cenário com métricas principais"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        cor = self._get_cor_cenario(dados.get('nivel_vies', 0))

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
        header.fill.solid()
        header.fill.fore_color.rgb = cor

        # Título do cenário
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = f"CENÁRIO {numero}\n{dados['titulo']}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Descrição
        desc = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.6))
        tf = desc.text_frame
        p = tf.paragraphs[0]
        p.text = dados['descricao']
        p.font.size = Pt(16)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

        # Cards de métricas (3 colunas)
        metricas = [
            ("Média Feminino", f"{dados['medias_depois']['Feminino']:.2f}", "Avaliação média mulheres"),
            ("Média Masculino", f"{dados['medias_depois']['Masculino']:.2f}", "Avaliação média homens"),
            ("Diferença", f"{abs(dados['diferenca_depois']):.3f}", "Gap entre gêneros")
        ]

        for i, (label, valor, desc_val) in enumerate(metricas):
            left = 0.5 + (i * 3.2)

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.7), Inches(3), Inches(1.8))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(245, 245, 245)

            tf = card.text_frame
            tf.clear()

            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(16)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = valor
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = cor
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = desc_val
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

        # Métricas adicionais (2 colunas)
        metricas2 = [
            ("P-value", f"{dados['p_value_depois']:.4f}", "Significância estatística"),
            ("Status", "✅ SEM VIÉS" if not dados['vies_detectado'] else "⚠️ VIÉS", "Resultado da análise")
        ]

        for i, (label, valor, desc_val) in enumerate(metricas2):
            left = 2 + (i * 3.2)

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.8), Inches(3), Inches(1.5))
            card.fill.solid()

            if label == "Status":
                if dados['vies_detectado']:
                    card.fill.fore_color.rgb = RGBColor(255, 235, 238)
                else:
                    card.fill.fore_color.rgb = RGBColor(232, 245, 233)
            else:
                card.fill.fore_color.rgb = RGBColor(245, 245, 245)

            tf = card.text_frame
            tf.clear()

            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = valor
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = desc_val
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER

        # Rodapé
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "💡 Próximos slides: 8 gráficos detalhados deste cenário"
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

        print(f"✓ Slide: Introdução {dados['titulo']}")

    def _slide_grafico(
        self,
        prs: Presentation,
        numero_cenario: int,
        titulo_grafico: str,
        caminho_imagem: Path,
        descricao: str
    ):
        """Slide com um gráfico"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Título pequeno
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = f"Cenário {numero_cenario} | {titulo_grafico}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Gráfico (grande)
        if caminho_imagem and caminho_imagem.exists():
            slide.shapes.add_picture(str(caminho_imagem), Inches(0.5), Inches(1.1), width=Inches(9))
        else:
            # Placeholder
            box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Gráfico não disponível:\n{caminho_imagem}"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

        # Descrição/análise
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📊 {descricao}"
        p.font.size = Pt(13)
        p.font.italic = True

        print(f"✓ Slide: {titulo_grafico}")

    def gerar_apresentacao_completa(
        self,
        dados_cenarios: Dict,
        todos_graficos: Dict,
        nome_arquivo: Optional[str] = None
    ) -> Path:
        """Gera apresentação PowerPoint ULTRA DETALHADA"""
        if nome_arquivo is None:
            nome_arquivo = f"apresentacao_completa_{self.timestamp}.pptx"

        caminho = self.output_dir / nome_arquivo

        print("\n=== Gerando Apresentação PowerPoint ULTRA DETALHADA ===\n")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Capa
        self._slide_capa(prs)

        # Slide 2: Índice
        self._slide_indice(prs, dados_cenarios)

        # Para cada cenário: intro + 8 gráficos
        for key in sorted(dados_cenarios.keys()):
            num = int(key.split('_')[1])
            dados = dados_cenarios[key]

            # Slide de introdução do cenário
            self._slide_intro_cenario(prs, num, dados)

            # 8 slides de gráficos
            graficos_cenario = todos_graficos.get(key, [])

            descricoes_graficos = [
                ("Distribuição de Scores (Antes)", "Mostra a distribuição antes de qualquer ajuste"),
                ("Distribuição de Scores (Depois)", "Distribuição após aplicação do nível de viés"),
                ("Comparação de Médias", "Comparação direta entre médias de homens e mulheres"),
                ("Boxplot por Tipo de Avaliação", "Distribuição de scores por diferentes tipos de avaliação"),
                ("Eficácia da Correção", "Análise da diferença de médias e significância estatística"),
                ("Histograma de Distribuição", "Distribuição geral com curva normal e percentis"),
                ("Desempenho vs Potencial", "Matriz Nine Box mostrando posicionamento"),
                ("Comparativo de Cenários", "Comparação de métricas chave")
            ]

            for idx, (titulo_graf, desc) in enumerate(descricoes_graficos):
                if idx < len(graficos_cenario):
                    caminho_img = graficos_cenario[idx]
                    self._slide_grafico(prs, num, titulo_graf, caminho_img, desc)

        # Salvar
        prs.save(str(caminho))

        print(f"\n✓ Apresentação COMPLETA salva: {caminho}")
        print(f"  Total de slides: {len(prs.slides)}")

        return caminho
